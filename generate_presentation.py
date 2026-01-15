#!/usr/bin/env python3
"""
Generate a PowerPoint presentation showcasing the dougboyd.com.au website build process
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme (matching the "Bright Meadow" theme)
    PRIMARY_COLOR = RGBColor(255, 107, 107)  # var(--color-primary)
    DARK_COLOR = RGBColor(43, 45, 66)        # var(--color-primary-dark)
    ACCENT_COLOR = RGBColor(255, 159, 28)    # var(--color-accent)
    TEXT_COLOR = RGBColor(51, 51, 51)        # var(--color-text)

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title box
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Building dougboyd.com.au"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Static Website Build Journey with Claude Code"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = TEXT_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Add date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = "January 2026"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = TEXT_COLOR
    date_para.alignment = PP_ALIGN.CENTER

    # Slide 2: Project Overview
    slide = add_title_content_slide(prs, "Project Overview", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Personal website showcasing multiple interest areas",
        "Pure static HTML/CSS/JavaScript - no framework dependencies",
        "Deployed to Azure Web Apps",
        "Content-driven with markdown source files",
        "Completely regenerated with each build",
        "AI-assisted content generation using Claude Code"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 3: Architecture Decisions
    slide = add_title_content_slide(prs, "Architecture Decisions", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Static site (no build process, no dependencies)",
        "Content separation (Page_Content/ directory structure)",
        "Tree-based navigation matching directory structure",
        "Azure Web Apps deployment with web.config",
        "Cache busting via query parameters (?v=YYYYMMDD)",
        "Section-specific background image organization"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 4: Key Technical Features
    slide = add_title_content_slide(prs, "Key Technical Features", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "CSS custom properties for theming",
        "Mobile-first responsive design",
        "Click-to-enlarge image modals on all content images",
        "Smooth scroll animations with Intersection Observer",
        "Background images with opacity overlays",
        "Google Analytics 4 tracking integration"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 5: Content Strategy
    slide = add_title_content_slide(prs, "Content Strategy", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "AI-assisted writing with directive markers:",
        "  • ***Creative Writing with AI - Doug's voice",
        "  • ***Formal Writing with AI - Technical content",
        "  • ***No AI - Raw, unmodified writing",
        "  • ***Instructional - Build instructions",
        "Monthly archive structure for time-based content",
        "Newsletter integration via Buttondown"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 6: Build Process - Phase 1
    slide = add_title_content_slide(prs, "Two-Phase Build Process", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Phase 1: Homepage Preview",
        "  • Generate new design system (colors, fonts, layout)",
        "  • Create homepage only for design approval",
        "  • No deployment - local preview only",
        "",
        "Phase 2: Full Site Build",
        "  • Apply approved design to all pages",
        "  • Regenerate all content from Page_Content/",
        "  • Deploy to Azure production",
        "  • Send email notifications to subscribers"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 7: Image Management
    slide = add_title_content_slide(prs, "Image Management Strategy", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Directory structure by section:",
        "  • images/background-images/fitness/",
        "  • images/background-images/rv7/",
        "  • images/background-images/generic/",
        "  • images/background-images/everything_else/",
        "",
        "CSS implementation via body classes",
        "Fallback to generic/ when section images unavailable",
        "All content images have click-to-enlarge modals"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 8: Modal Image Feature
    slide = add_title_content_slide(prs, "Image Modal Implementation", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "User Experience:",
        "  • Click any image to view full-size",
        "  • Multiple close methods (ESC, click outside, × button)",
        "  • Prevents background scrolling when open",
        "  • Smooth animations (fadeIn, zoomIn)",
        "",
        "Technical Implementation:",
        "  • CSS ::before pseudo-element for overlay",
        "  • JavaScript event handlers for interactions",
        "  • Hover effects indicate clickability"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 9: Footer Enhancement
    slide = add_title_content_slide(prs, "Footer Design Pattern", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Required footer elements on every page:",
        "  • Copyright notice",
        "  • Subscribe link (to #newsletter anchor)",
        "  • Ko-fi support button",
        "  • Generation timestamp",
        "",
        "Subscribe link navigation:",
        "  • Homepage: href='#newsletter'",
        "  • Other pages: href='index.html#newsletter'",
        "  • Subdirectories: relative paths",
        "",
        "Background image at 8% opacity for visual interest"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 10: Fitness Section Restructure
    slide = add_title_content_slide(prs, "Fitness Section Evolution", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Initial State:",
        "  • Single page with training plan and progress log",
        "",
        "Restructured Design:",
        "  • fitness.html - Overview and training plan",
        "  • fitness-january-2026.html - Monthly archive",
        "  • Scalable structure for future months",
        "",
        "Benefits:",
        "  • Cleaner separation of concerns",
        "  • Better organization for long-term tracking",
        "  • Easier navigation to specific time periods"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 11: Accountability Feature
    slide = add_title_content_slide(prs, "Accountability Log Addition", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "New 'Everything Else' section launched",
        "Daily accountability log tracking:",
        "  • Activities and progress",
        "  • Honest self-assessment",
        "  • Personal narrative",
        "",
        "Image layout strategy:",
        "  • Floating images (45% width, max 400px)",
        "  • Alternating left/right placement",
        "  • Text wrapping around images",
        "  • Click-to-enlarge modal functionality"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 12: Deployment Pipeline
    slide = add_title_content_slide(prs, "Deployment Pipeline", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "scripts/deploy-azure.sh",
        "  • Creates deployment artifact",
        "  • Deploys to Azure Web App via zip",
        "  • Commits changes to git repository",
        "  • Pushes to GitHub",
        "",
        "scripts/send-notification.sh",
        "  • Standalone email notification system",
        "  • 8 random email template variations",
        "  • Sends to all Buttondown subscribers",
        "",
        "Combined: bash scripts/deploy-azure.sh prod && bash scripts/send-notification.sh"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 13: CSS Architecture
    slide = add_title_content_slide(prs, "CSS Architecture", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Single main.css file with:",
        "  • CSS custom properties in :root",
        "  • Mobile-first responsive design",
        "  • Utility classes for spacing",
        "  • Smooth transitions throughout",
        "",
        "Background image implementation:",
        "  • Default .hero-background class",
        "  • Body class overrides (specificity)",
        "  • Example: body.fitness-page .hero-background",
        "",
        "Footer background via ::before pseudo-element"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 14: JavaScript Features
    slide = add_title_content_slide(prs, "JavaScript Features", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "No external dependencies - vanilla JavaScript",
        "",
        "Key functionality:",
        "  • Mobile navigation toggle",
        "  • Smooth scrolling for anchor links",
        "  • Intersection Observer for fade-in animations",
        "  • Image modal system (click-to-enlarge)",
        "  • Background scroll prevention when modal active",
        "",
        "Performance considerations:",
        "  • Event delegation where appropriate",
        "  • Minimal DOM manipulation"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 15: Cache Busting Strategy
    slide = add_title_content_slide(prs, "Cache Busting Strategy", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Problem:",
        "  • Browsers aggressively cache CSS/JS files",
        "  • Users see outdated styles after deployment",
        "",
        "Solution:",
        "  • Query parameters with build date",
        "  • Example: main.css?v=20260112",
        "  • Format: YYYYMMDD",
        "",
        "Implementation:",
        "  • All <link> tags include version parameter",
        "  • All <script> tags include version parameter",
        "  • Version updated with each build"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 16: Background Image Technical Details
    slide = add_title_content_slide(prs, "Background Image Technical Details", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Hero sections:",
        "  • 15% opacity for readability",
        "  • background-size: cover",
        "  • background-position: center",
        "  • Separate <div class='hero-background'> element",
        "",
        "Footer:",
        "  • 8% opacity (darker base color)",
        "  • ::before pseudo-element implementation",
        "  • z-index layering for content above image",
        "",
        "CSS specificity ensures correct section images load"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 17: Content Generation Workflow
    slide = add_title_content_slide(prs, "Content Generation Workflow", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "1. Content creation in Page_Content/ directory",
        "   • Markdown files with AI directives",
        "   • Organized by section subdirectories",
        "",
        "2. AI processing during build:",
        "   • Creative writing in Doug's voice",
        "   • Formal/technical content generation",
        "   • Preserve raw writing where marked",
        "",
        "3. HTML generation:",
        "   • Apply current design system",
        "   • Insert Google Analytics tracking",
        "   • Add image modal functionality",
        "   • Generate footer with timestamp"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 18: Design Philosophy
    slide = add_title_content_slide(prs, "Design Philosophy", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Current theme: 'Bright Meadow'",
        "  • Lighter, more fun aesthetic",
        "  • Bright, playful, energetic, approachable",
        "  • Avoid dark/brutalist designs",
        "",
        "Inspiration sources:",
        "  • arlenmccluskey.com",
        "  • mackandpouya.com",
        "  • aileen.co",
        "  • iamtamara.design",
        "",
        "Every build = completely new design",
        "Colors, typography, layout patterns all regenerated"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 19: Security Considerations
    slide = add_title_content_slide(prs, "Security Considerations", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Never commit sensitive files:",
        "  • .env (contains API keys)",
        "  • .godaddy-config",
        "  • Credentials or secrets",
        "",
        "Azure web.config security headers:",
        "  • X-Frame-Options",
        "  • X-Content-Type-Options",
        "  • HTTPS redirection",
        "",
        "Static compression enabled",
        "Custom 404 error page configured"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 20: Analytics & Tracking
    slide = add_title_content_slide(prs, "Analytics & Tracking", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Google Analytics 4 integration:",
        "  • Tracking ID: G-X8VGBX1GHX",
        "  • Snippet in <head> of every page",
        "  • Tracks visitor counts and page views",
        "",
        "Newsletter tracking:",
        "  • Buttondown integration for subscribers",
        "  • Email notifications on new builds",
        "  • 8 randomized email templates",
        "",
        "Ko-fi support button:",
        "  • https://ko-fi.com/dougboyd",
        "  • Present on every page footer"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 21: Key Challenges & Solutions
    slide = add_title_content_slide(prs, "Key Challenges & Solutions", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Challenge: Images not showing after build",
        "Solution: Implemented body class CSS specificity pattern",
        "",
        "Challenge: Browser caching old CSS/JS",
        "Solution: Query parameter cache busting with build date",
        "",
        "Challenge: Images full-width with no text flow",
        "Solution: Floating layout (45% width, alternating sides)",
        "",
        "Challenge: Making images feel interactive",
        "Solution: Click-to-enlarge modal with multiple close methods",
        "",
        "Challenge: Footer needed Subscribe link",
        "Solution: Anchor navigation to #newsletter section"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 22: Build Documentation
    slide = add_title_content_slide(prs, "Build Documentation", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "buildInstructions.md:",
        "  • Two-phase build process documented",
        "  • Background image requirements",
        "  • Image modal treatment specification",
        "  • Footer requirements and structure",
        "  • Cache busting requirements",
        "  • Google Analytics integration",
        "",
        "CLAUDE.md:",
        "  • Project overview for Claude Code",
        "  • Architecture decisions",
        "  • Development workflow",
        "  • Deployment instructions"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 23: Testing & Validation
    slide = add_title_content_slide(prs, "Testing & Validation", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Local development:",
        "  • python -m http.server 8000",
        "  • npx http-server",
        "  • No build step required",
        "",
        "Validation checks:",
        "  • All navigation links functional",
        "  • Background images loading correctly",
        "  • Image modals working on all pages",
        "  • Subscribe links pointing to #newsletter",
        "  • Responsive design on mobile/tablet/desktop",
        "  • Cache busting parameters present",
        "",
        "Deployment validation:",
        "  • Azure status check via scripts/deploy-azure.sh status"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 24: Future Enhancements
    slide = add_title_content_slide(prs, "Future Enhancements", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Planned additions:",
        "  • RV7 builder's log detailed entries",
        "  • More writings and technical articles",
        "  • Additional 'Everything Else' projects",
        "  • Family planner development notes",
        "  • Software experiments documentation",
        "",
        "Infrastructure improvements:",
        "  • Automated testing pipeline",
        "  • Performance monitoring",
        "  • A/B testing for design variations",
        "",
        "Content strategy:",
        "  • Regular running log updates",
        "  • Daily accountability entries"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 25: Lessons Learned
    slide = add_title_content_slide(prs, "Lessons Learned", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Simplicity wins:",
        "  • No framework = no framework problems",
        "  • Static files = fast, reliable, easy to deploy",
        "",
        "Plan for iteration:",
        "  • Two-phase build allows design exploration",
        "  • Separation of content and presentation",
        "",
        "Document everything:",
        "  • Build instructions prevent rework",
        "  • Clear patterns enable consistency",
        "",
        "AI as a force multiplier:",
        "  • Content generation with voice preservation",
        "  • Rapid prototyping of features",
        "  • Consistent implementation patterns"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 26: Technical Stack Summary
    slide = add_title_content_slide(prs, "Technical Stack Summary", PRIMARY_COLOR, TEXT_COLOR)

    content = [
        "Frontend:",
        "  • Pure HTML5, CSS3, JavaScript (ES6+)",
        "  • No frameworks, no dependencies",
        "",
        "Hosting & Deployment:",
        "  • Azure Web Apps (Production)",
        "  • GitHub for version control",
        "  • Bash scripts for deployment automation",
        "",
        "Third-party Services:",
        "  • Google Analytics 4 (tracking)",
        "  • Buttondown (newsletter)",
        "  • Ko-fi (support/donations)",
        "",
        "Development Tools:",
        "  • Claude Code (AI-assisted development)",
        "  • Python for automation scripts"
    ]
    add_bullet_points(slide, content, TEXT_COLOR)

    # Slide 27: Closing Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add closing message
    closing_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    closing_frame = closing_box.text_frame
    closing_frame.text = "Questions?"
    closing_para = closing_frame.paragraphs[0]
    closing_para.font.size = Pt(54)
    closing_para.font.bold = True
    closing_para.font.color.rgb = PRIMARY_COLOR
    closing_para.alignment = PP_ALIGN.CENTER

    # Add contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2))
    contact_frame = contact_box.text_frame

    p = contact_frame.paragraphs[0]
    p.text = "dougboyd.com.au"
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    p = contact_frame.add_paragraph()
    p.text = "github.com/dougboyd/dougboyd.com.au"
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

    p = contact_frame.add_paragraph()
    p.text = ""

    p = contact_frame.add_paragraph()
    p.text = "Built with Claude Code"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER

    return prs

def add_title_content_slide(prs, title_text, title_color, text_color):
    """Add a slide with title and content area"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = title_color

    return slide

def add_bullet_points(slide, bullet_points, text_color):
    """Add bullet points to a slide"""
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # Determine indentation level based on leading spaces
        indent_level = 0
        stripped_point = point.lstrip()
        if point.startswith('  •'):
            indent_level = 1
            stripped_point = stripped_point[2:].strip()
        elif point.startswith('•'):
            stripped_point = stripped_point[1:].strip()

        p.text = stripped_point
        p.level = indent_level
        p.font.size = Pt(16) if indent_level == 0 else Pt(14)
        p.font.color.rgb = text_color

        if point.strip() == "":
            p.text = " "  # Empty line for spacing

if __name__ == "__main__":
    print("Generating PowerPoint presentation...")

    try:
        prs = create_presentation()
        output_file = "dougboyd_website_build.pptx"
        prs.save(output_file)
        print(f"✓ Presentation saved as: {output_file}")
        print(f"  Total slides: {len(prs.slides)}")
    except Exception as e:
        print(f"✗ Error generating presentation: {e}")
        import traceback
        traceback.print_exc()
